"""
Enhanced Analysis Overview - Streamlit Cloud Ready
ABC + Daily Volumes + Bin Estimator + Order Batching + PowerPoint Report
"""

import io
import csv
import math
import statistics
from datetime import datetime
from typing import Optional, Tuple, List, Dict, Any
from collections import defaultdict

import numpy as np
import pandas as pd
import streamlit as st
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from pptx import Presentation
from pptx.util import Inches, Pt
from datasketch import MinHashLSH, MinHash

from abc_core import run_abc, run_lines_per_day


# ============================================================================
# CONFIGURATION
# ============================================================================

class Config:
    PAGE_TITLE = "Analysis Overview"
    PAGE_LAYOUT = "wide"
    PPT_WIDTH, PPT_HEIGHT = 13.333, 7.5
    
    COLORS = {
        'primary': '#EB0A1E', 'dark': '#111111', 'grey': '#6B7280',
        'grid': '#E5E7EB', 'background': '#FFFFFF', 'accent': '#1D4ED8',
    }
    
    DEFAULT_A_CUTOFF, DEFAULT_B_CUTOFF = 0.80, 0.95
    DEFAULT_PERCENTILES = [0.50, 0.80, 1.00]
    DEFAULT_GROWTH = 0.0
    DEFAULT_BIN_UTIL, DEFAULT_ITEMS_PER_BIN = 0.60, 50.0
    DEFAULT_STOCK_ON_HAND_DAYS, DEFAULT_SAFETY_STOCK = 14, 0.10
    DEFAULT_PLANNING_BASIS = "P95"
    DEFAULT_BIN_L, DEFAULT_BIN_W, DEFAULT_BIN_H = 600.0, 400.0, 330.0
    DEFAULT_WAVING_TIMER, DEFAULT_GROUP_SIZE = 0, 4
    DEFAULT_THRESHOLD, DEFAULT_PERMUTATIONS = 0.1, 128
    
    CSV_ENCODING_OPTIONS = ["Auto", "utf-8", "utf-8-sig", "cp1252", "latin-1"]
    CSV_DELIMITER_OPTIONS = ["Auto", ",", ";", "\t", "|"]
    CSV_MAX_SNIFF_BYTES = 200_000


plt.rcParams.update({'font.family': 'DejaVu Sans', 'axes.titlesize': 11, 'axes.titleweight': 'bold'})


# ============================================================================
# UTILITIES
# ============================================================================

def format_integer(value): 
    try: return f"{int(round(float(value))):,}"
    except: return "N/A"

def format_number(value, decimals=0): 
    try: return f"{float(value):,.{decimals}f}"
    except: return "N/A"

def format_percentage(value): 
    try: return f"{float(value):.0%}"
    except: return "N/A"

def safe_to_datetime_date(series):
    return pd.to_datetime(series, errors="coerce").dt.date

def first_existing_col(df, candidates):
    cols = {c.lower(): c for c in df.columns}
    for cand in candidates:
        if cand.lower() in cols:
            return cols[cand.lower()]
    return None

def clamp01(x): 
    return max(0.0, min(1.0, float(x)))

def df_to_csv_bytes(df):
    buf = io.StringIO()
    df.to_csv(buf, index=False)
    return buf.getvalue().encode("utf-8")

def coerce_to_yyyy_mm_dd(series):
    """Force daily date labels to YYYY-MM-DD format."""
    s = series.copy()
    if pd.api.types.is_datetime64_any_dtype(s):
        dt = s
    else:
        s_str = s.astype(str).str.strip()
        dt1 = pd.to_datetime(s_str, errors="coerce", format="%Y-%m-%d")
        dt2 = pd.to_datetime(s_str, errors="coerce", format="%Y-%d-%m")
        dt = dt1 if dt1.isna().sum() <= dt2.isna().sum() else dt2
        if dt.isna().all():
            dt = pd.to_datetime(s_str, errors="coerce")
    return dt.dt.strftime("%Y-%m-%d")


# ============================================================================
# BATCHING FUNCTIONS
# ============================================================================

def round_up_datetime(dt, interval):
    total_minutes = dt.hour * 60 + dt.minute
    rounded_minutes = ((total_minutes + interval - 1) // interval) * interval % (24 * 60)
    return dt.replace(hour=rounded_minutes // 60, minute=rounded_minutes % 60, second=0, microsecond=0)

def get_minhash(item_set, num_perm):
    m = MinHash(num_perm=num_perm)
    for item in item_set:
        m.update(str(item).encode('utf8'))
    return m

def _exact_jaccard(a, b):
    if not a and not b: return 0.0
    u = len(a | b)
    return len(a & b) / u if u else 0.0

def _merge_small_groups_exact(groups, client_items, GroupSize, threshold):
    full_groups = [g for g in groups if len(g) == GroupSize]
    small_ids = [i for i, g in enumerate(groups) if len(g) < GroupSize]
    if len(small_ids) <= 1: return groups
    
    union_sets = {gid: set().union(*(client_items[c] for c in groups[gid])) for gid in small_ids}
    sizes = {gid: len(groups[gid]) for gid in small_ids}
    
    neighbors = {gid: {} for gid in small_ids}
    for i, a in enumerate(small_ids):
        for b in small_ids[i + 1:]:
            jac = _exact_jaccard(union_sets[a], union_sets[b])
            if jac >= threshold and sizes[a] + sizes[b] <= GroupSize:
                neighbors[a][b] = neighbors[b][a] = jac
    
    unmerged, merged_small = set(small_ids), []
    while unmerged:
        seed = max(unmerged, key=lambda x: (len(neighbors[x]), max(neighbors[x].values()) if neighbors[x] else 0))
        bloc = [seed]
        total_size = sizes[seed]
        
        for n in sorted([n for n in neighbors[seed] if n in unmerged], key=lambda n: neighbors[seed][n], reverse=True):
            if total_size + sizes[n] <= GroupSize:
                bloc.append(n)
                total_size += sizes[n]
        
        for gid in bloc: unmerged.discard(gid)
        merged_small.append([c for gid in bloc for c in groups[gid]])
    
    return full_groups + merged_small

def _pack_leftovers_fill_whatever(groups, GroupSize):
    full = [g for g in groups if len(g) == GroupSize]
    rest = [c for g in groups if len(g) < GroupSize for c in g]
    return full + [rest[i:i + GroupSize] for i in range(0, len(rest), GroupSize)]

def group_clients_by_similarity(client_items, GroupSize, threshold, permutations):
    minhashes = {client: get_minhash(items, permutations) for client, items in client_items.items()}
    lsh = MinHashLSH(threshold=threshold, num_perm=permutations)
    for client, m in minhashes.items():
        lsh.insert(client, m)
    
    similarity_graph = defaultdict(set)
    for client, m in minhashes.items():
        for c in lsh.query(m):
            if c != client:
                similarity_graph[client].add(c)
    
    degree = {client: len(neighbors) for client, neighbors in similarity_graph.items()}
    ungrouped, groups = set(client_items.keys()), []
    
    while len(ungrouped) >= GroupSize:
        c = max(ungrouped, key=lambda x: degree.get(x, 0))
        neigh = list(similarity_graph[c] & ungrouped)
        group = [c] + neigh[:GroupSize - 1]
        need = GroupSize - len(group)
        if need > 0:
            group += list(ungrouped - set(group))[:need]
        groups.append(group)
        ungrouped -= set(group)
    
    if ungrouped: groups.append(list(ungrouped))
    groups = _merge_small_groups_exact(groups, client_items, GroupSize, threshold)
    return _pack_leftovers_fill_whatever(groups, GroupSize)

def group_clients_by_day_optimized(df, GroupSize, num_perm, threshold, order_column_name, sku_column_name):
    results = []
    for dt_val in sorted(df["roundeddatetime"].unique()):
        group = df[df["roundeddatetime"] == dt_val]
        client_items = group.groupby(order_column_name)[sku_column_name].apply(set).to_dict()
        
        for group_id, grp in enumerate(group_clients_by_similarity(client_items, GroupSize, threshold, num_perm), 1):
            unique_items, total_items = set(), 0
            for c in grp:
                items = client_items.get(c, set())
                if isinstance(items, set):
                    total_items += len(items)
                    unique_items |= items
            
            results.append({
                "roundeddatetime": dt_val, "Taskgroup": group_id, "Order": grp,
                "Unique Sku total in order": len(unique_items), "total_items": total_items,
                "unique_items": unique_items, "Batch Factor": total_items / len(unique_items) if unique_items else 0,
            })
    return results


# ============================================================================
# CHART STYLING
# ============================================================================

class ChartStyler:
    @staticmethod
    def apply_style(ax):
        ax.set_facecolor(Config.COLORS['background'])
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        for spine in ['left', 'bottom']:
            ax.spines[spine].set_color(Config.COLORS['grid'])
            ax.spines[spine].set_linewidth(1.0)
        ax.grid(True, axis='y', color=Config.COLORS['grid'], linewidth=0.8)
        ax.grid(False, axis='x')
        ax.tick_params(axis='both', colors=Config.COLORS['grey'], labelsize=9)
        for label in [ax.xaxis.label, ax.yaxis.label]:
            label.set_color(Config.COLORS['dark'])
            label.set_fontsize(10)

    @staticmethod
    def add_legend(ax, loc='upper left'):
        legend = ax.legend(frameon=True, fancybox=True, framealpha=0.95, fontsize=9, 
                          loc=loc, borderpad=0.6, labelspacing=0.5, handlelength=2.2)
        if legend:
            legend.get_frame().set_edgecolor(Config.COLORS['grid'])
            legend.get_frame().set_linewidth(1.0)
            legend.get_frame().set_facecolor(Config.COLORS['background'])
            for text in legend.get_texts():
                text.set_color(Config.COLORS['dark'])


# ============================================================================
# DATA LOADING
# ============================================================================

class DataLoader:
    @staticmethod
    def detect_text_encoding(file_bytes):
        if file_bytes.startswith(b'\xef\xbb\xbf'): return "utf-8-sig"
        for enc in ("utf-8", "cp1252"):
            try: 
                file_bytes.decode(enc)
                return enc
            except: pass
        return "latin-1"

    @staticmethod
    def detect_csv_delimiter(sample):
        try: return csv.Sniffer().sniff(sample, delimiters=[',', ';', '\t', '|']).delimiter
        except: return ','

    @staticmethod
    def _read_csv_multi_pass(file_bytes, delimiter, encoding):
        kwargs = dict(sep=delimiter, encoding=encoding)
        for engine, extra in [("c", {}), ("python", {}), ("python", {"on_bad_lines": "skip"})]:
            try: return pd.read_csv(io.BytesIO(file_bytes), engine=engine, **kwargs, **extra)
            except: pass
        try: return pd.read_csv(io.BytesIO(file_bytes), engine="python", on_bad_lines="skip", 
                               encoding=encoding, encoding_errors="replace", sep=delimiter)
        except:
            text = file_bytes.decode(encoding, errors="replace")
            return pd.read_csv(io.StringIO(text), engine="python", sep=delimiter, on_bad_lines="skip")

    @staticmethod
    def load_csv(file_bytes, filename, encoding_override=None, delimiter_override=None):
        encoding = encoding_override if (encoding_override and encoding_override != "Auto") else DataLoader.detect_text_encoding(file_bytes)
        if delimiter_override and delimiter_override != "Auto":
            delimiter = delimiter_override
        else:
            sample = file_bytes[:Config.CSV_MAX_SNIFF_BYTES].decode(encoding, errors="replace")
            delimiter = DataLoader.detect_csv_delimiter(sample)
        
        df = DataLoader._read_csv_multi_pass(file_bytes, delimiter, encoding)
        return df, f"Loaded CSV: {filename} (delimiter={repr(delimiter)}, encoding='{encoding}')"

    @staticmethod
    def load_excel(file_obj, sheet_name):
        return pd.read_excel(file_obj, sheet_name=sheet_name)


# ============================================================================
# CHART GENERATORS
# ============================================================================

class ChartGenerator:
    @staticmethod
    def create_abc_chart(abc_df, abc_summary, a_cutoff, b_cutoff, figsize=(7.2, 3.3), dpi=140):
        abc_plot = abc_df.reset_index(drop=True).copy()
        
        if 'SKU_Pct' in abc_plot.columns:
            x = abc_plot['SKU_Pct'].to_numpy() * 100.0
        else:
            n = len(abc_plot)
            total_skus = max(int(abc_summary.get('total_skus', n)), n)
            x = (np.arange(1, n + 1) / float(total_skus)) * 100.0
        
        y = abc_plot['cum_share'].to_numpy() * 100.0
        x, y = np.concatenate(([0.0], x)), np.concatenate(([0.0], y))
        
        cut_a_pct, cut_b_pct = a_cutoff * 100.0, b_cutoff * 100.0
        idx_a = int(np.argmax(y >= cut_a_pct)) if np.any(y >= cut_a_pct) else len(y) - 1
        idx_b = int(np.argmax(y >= cut_b_pct)) if np.any(y >= cut_b_pct) else len(y) - 1
        x_a, x_b = float(x[idx_a]), float(x[idx_b])
        
        fig, ax = plt.subplots(figsize=figsize, dpi=dpi)
        ax.plot(x, y, linewidth=2.2, color=Config.COLORS['dark'], label='Cumulative % Lines')
        
        ax.hlines(cut_a_pct, 0, x_a, linestyles='--', linewidth=1.8, colors=Config.COLORS['primary'], label=f"A cutoff ({int(cut_a_pct)}%)")
        ax.vlines(x_a, 0, cut_a_pct, linestyles='--', linewidth=1.8, colors=Config.COLORS['primary'])
        ax.scatter([x_a], [cut_a_pct], s=26, color=Config.COLORS['primary'], zorder=5)
        
        ax.hlines(cut_b_pct, 0, x_b, linestyles='--', linewidth=1.8, colors=Config.COLORS['primary'], alpha=0.70, label=f"B cutoff ({int(cut_b_pct)}%)")
        ax.vlines(x_b, 0, cut_b_pct, linestyles='--', linewidth=1.8, colors=Config.COLORS['primary'], alpha=0.70)
        ax.scatter([x_b], [cut_b_pct], s=26, color=Config.COLORS['primary'], alpha=0.70, zorder=5)
        
        ax.set_xlim(0, 100)
        ax.set_ylim(0, 100)
        ax.set_xlabel('SKU (%)')
        ax.set_ylabel('% Lines (cumulative)')
        ax.set_xticks(range(0, 101, 10))
        ax.set_yticks(range(0, 101, 10))
        
        ChartStyler.apply_style(ax)
        ChartStyler.add_legend(ax, loc='lower right')
        fig.tight_layout(pad=1.0)
        return fig, x_a, x_b

    @staticmethod
    def create_daily_volume_chart(daily_df, y_column, y_label, percentiles, growth_rate, figsize=(7.2, 3.3), dpi=140):
        plot_df = daily_df[['Date', y_column]].copy()
        
        if plot_df['Date'].dtype != "datetime64[ns]":
            plot_df['Date'] = pd.to_datetime(plot_df['Date'], errors='coerce', format="%Y-%m-%d")
            if plot_df['Date'].isna().all():
                plot_df['Date'] = pd.to_datetime(plot_df['Date'], errors='coerce')
        
        plot_df[y_column] = pd.to_numeric(plot_df[y_column], errors='coerce')
        plot_df = plot_df.dropna(subset=['Date', y_column]).sort_values('Date')
        
        growth_factor = 1.0 + (growth_rate / 100.0)
        plot_df['Projected'] = plot_df[y_column] * growth_factor
        
        fig, ax = plt.subplots(figsize=figsize, dpi=dpi)
        ax.plot(plot_df['Date'], plot_df[y_column], linewidth=2.2, color=Config.COLORS['dark'], label=f"Actual {y_label}")
        
        if growth_rate != 0.0:
            ax.plot(plot_df['Date'], plot_df['Projected'], linewidth=2.2, 
                   color=Config.COLORS['accent'], alpha=0.95, label=f"Projected ({growth_rate:+.0f}%)")
        
        ax.xaxis.set_major_formatter(mdates.DateFormatter('%Y-%m-%d'))
        fig.autofmt_xdate()
        
        basis_column = 'Projected' if growth_rate != 0.0 else y_column
        values = plot_df[basis_column].to_numpy()
        values = values[np.isfinite(values)]
        
        if values.size > 0:
            max_pct = max(percentiles) if percentiles else 1.0
            for pct in sorted(percentiles):
                y_line = float(np.quantile(values, pct))
                is_max = (pct == max_pct)
                ax.axhline(y_line, linestyle='--', linewidth=1.5 if is_max else 1.2,
                          color=Config.COLORS['primary'], alpha=0.85 if is_max else 0.55, label=f"P{int(pct * 100)}")
        
        ax.set_xlabel('Date')
        ax.set_ylabel(y_label)
        ChartStyler.apply_style(ax)
        ChartStyler.add_legend(ax, loc='upper left')
        fig.tight_layout(pad=1.0)
        return fig

    @staticmethod
    def create_batch_factor_chart(batch_results, figsize=(7.2, 3.3), dpi=140):
        df_report = pd.DataFrame(batch_results)
        df_report = df_report[df_report["Taskgroup"].notna()].copy()
        
        dates = sorted(df_report["roundeddatetime"].unique())
        avg_batch_factors = []
        for date in dates:
            batch_factors = df_report[df_report["roundeddatetime"] == date]["Batch Factor"].tolist()
            avg_batch_factors.append(sum(batch_factors) / len(batch_factors) if batch_factors else 0)
        
        overall_median = statistics.median(avg_batch_factors) if avg_batch_factors else 0
        
        fig, ax = plt.subplots(figsize=figsize, dpi=dpi)
        ax.scatter(dates, avg_batch_factors, s=50, color=Config.COLORS['primary'], alpha=0.7)
        ax.axhline(overall_median, color=Config.COLORS['accent'], linestyle='--', 
                  linewidth=1.5, label=f'Overall Median ({overall_median:.4f})')
        
        ax.set_xlabel('Date')
        ax.set_ylabel('Average Batch Factor')
        ax.set_title('Average Batch Factor Per Day')
        ChartStyler.apply_style(ax)
        ChartStyler.add_legend(ax, loc='upper left')
        fig.autofmt_xdate()
        fig.tight_layout(pad=1.0)
        return fig


# ============================================================================
# POWERPOINT GENERATOR
# ============================================================================

class PowerPointGenerator:
    @staticmethod
    def figure_to_bytes(fig):
        buf = io.BytesIO()
        fig.savefig(buf, format='png', bbox_inches='tight', dpi=220)
        buf.seek(0)
        return buf.getvalue()

    @staticmethod
    def add_slide_title(slide, title):
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.2), Inches(0.7))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size, p.font.bold = Pt(30), True

    @staticmethod
    def add_two_column_slide(prs, title, bullets, image_bytes=None):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        PowerPointGenerator.add_slide_title(slide, title)
        
        bullet_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(6.1), Inches(5.8))
        tf = bullet_box.text_frame
        tf.clear()
        
        heading = tf.paragraphs[0]
        heading.text = "Key Points"
        heading.font.size, heading.font.bold = Pt(18), True
        heading.space_after = Pt(8)
        
        for b in bullets:
            para = tf.add_paragraph()
            para.text = b
            para.level, para.font.size = 0, Pt(14)
            para.space_after = Pt(4)
        
        if image_bytes:
            slide.shapes.add_picture(io.BytesIO(image_bytes), Inches(7.1), Inches(1.35), 
                                    width=Inches(5.8), height=Inches(5.8))

    @staticmethod
    def _summarise_data_window(df, date_col):
        try:
            dt = pd.to_datetime(df[date_col], errors='coerce').dropna()
            if dt.empty: return "N/A", "N/A", None
            dmin, dmax = dt.min().date().isoformat(), dt.max().date().isoformat()
            ndays = int((dt.max().normalize() - dt.min().normalize()).days) + 1
            return dmin, dmax, ndays
        except: return "N/A", "N/A", None

    @staticmethod
    def _daily_stats_bullets(daily_df, y_col, y_label, percentiles, growth_rate):
        if daily_df is None or daily_df.empty or y_col not in daily_df.columns:
            return [f"{y_label}: no data available."]
        
        plot_df = daily_df[['Date', y_col]].copy()
        plot_df['Date'] = pd.to_datetime(plot_df['Date'], errors='coerce', format="%Y-%m-%d")
        if plot_df['Date'].isna().all():
            plot_df['Date'] = pd.to_datetime(plot_df['Date'], errors='coerce')
        plot_df[y_col] = pd.to_numeric(plot_df[y_col], errors='coerce')
        plot_df = plot_df.dropna(subset=['Date', y_col]).sort_values('Date')
        if plot_df.empty: return [f"{y_label}: no valid dates/values."]
        
        growth_factor = 1.0 + (growth_rate / 100.0)
        basis_col = 'Projected' if growth_rate != 0.0 else y_col
        if basis_col == 'Projected':
            plot_df['Projected'] = plot_df[y_col] * growth_factor
        
        peak_row = plot_df.loc[plot_df[basis_col].idxmax()]
        peak_value, peak_date = float(peak_row[basis_col]), peak_row['Date'].date().isoformat()
        
        vals = plot_df[basis_col].to_numpy()
        vals = vals[np.isfinite(vals)]
        p_lines = [f"P{int(pct*100)}: {format_integer(np.quantile(vals, pct))}" 
                  for pct in sorted(percentiles) if vals.size]
        
        basis_txt = f"Projected series (growth {growth_rate:+.0f}%)" if growth_rate != 0.0 else "Actual series (no growth)"
        bullets = [f"Basis: {basis_txt}.", f"Design day (peak): {peak_date} at {format_integer(peak_value)} {y_label}."]
        bullets.extend(p_lines[:6])
        if len(p_lines) > 6:
            bullets.append("Additional percentiles omitted for readability.")
        return bullets

    @staticmethod
    def generate_report(uploaded_name, sheet, sku_col, line_col, date_col, qty_col, a_cutoff, b_cutoff,
                       percentiles, growth_rate, abc_df, abc_summary, daily_df, raw_df, bin_summary=None,
                       bin_df=None, batch_summary=None, batch_results=None):
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(Config.PPT_WIDTH), Inches(Config.PPT_HEIGHT)
        
        dmin, dmax, ndays = PowerPointGenerator._summarise_data_window(raw_df, date_col)
        total_rows, total_cols = len(raw_df) if raw_df is not None else 0, len(raw_df.columns) if raw_df is not None else 0
        total_skus = abc_summary.get('total_skus') if isinstance(abc_summary, dict) else None
        pct_txt = ", ".join([f"P{int(p*100)}" for p in sorted(percentiles)]) if percentiles else "None"
        
        assumptions = [
            f"Source file: {uploaded_name}", f"Sheet: {sheet if sheet else '(n/a)'}",
            f"Rows analysed: {format_integer(total_rows)} | Columns: {format_integer(total_cols)}",
            f"Date column: '{date_col}' | Window: {dmin} to {dmax}" + (f" ({ndays} days)" if ndays else ""),
            f"SKU column: '{sku_col}' | Line column: '{line_col}'",
            f"Quantity column: '{qty_col}'" if qty_col else "Quantity column: (not provided)",
            f"Daily chart percentiles: {pct_txt}", f"Growth projection: {growth_rate:+.0f}%",
            "Design day definition: peak day of selected basis series.",
            f"ABC thresholds: A={a_cutoff:.0%}, B={b_cutoff:.0%}",
            f"Total SKUs (ABC): {format_integer(total_skus)}" if total_skus is not None else "Total SKUs (ABC): N/A",
        ]
        PowerPointGenerator.add_two_column_slide(prs, "Data Assumptions", assumptions)
        
        abc_fig, x_a, x_b = ChartGenerator.create_abc_chart(abc_df, abc_summary, a_cutoff, b_cutoff, figsize=(5.2, 5.2), dpi=170)
        abc_bytes = PowerPointGenerator.figure_to_bytes(abc_fig)
        plt.close(abc_fig)
        
        abc_bullets = [
            f"Total SKUs: {format_integer(abc_summary.get('total_skus'))}",
            f"Total lines: {format_number(abc_summary.get('total_value'), 0)}",
            f"A cutoff ({a_cutoff:.0%}) reached at ~{x_a:.1f}% of SKUs.",
            f"B cutoff ({b_cutoff:.0%}) reached at ~{x_b:.1f}% of SKUs.",
            "Interpretation: left = high movers; right = long tail SKUs."
        ]
        PowerPointGenerator.add_two_column_slide(prs, "ABC Distribution", abc_bullets, abc_bytes)
        
        lines_fig = ChartGenerator.create_daily_volume_chart(daily_df, 'Lines', 'Lines/day', percentiles, growth_rate, figsize=(5.2, 5.2), dpi=170)
        lines_bytes = PowerPointGenerator.figure_to_bytes(lines_fig)
        plt.close(lines_fig)
        lines_bullets = PowerPointGenerator._daily_stats_bullets(daily_df, 'Lines', 'lines/day', percentiles, growth_rate)
        PowerPointGenerator.add_two_column_slide(prs, "Lines per Day", lines_bullets, lines_bytes)
        
        has_quantity = 'Quantity' in daily_df.columns and daily_df['Quantity'].notna().any()
        if has_quantity:
            items_fig = ChartGenerator.create_daily_volume_chart(daily_df, 'Quantity', 'Items/day', percentiles, growth_rate, figsize=(5.2, 5.2), dpi=170)
            items_bytes = PowerPointGenerator.figure_to_bytes(items_fig)
            plt.close(items_fig)
            items_bullets = PowerPointGenerator._daily_stats_bullets(daily_df, 'Quantity', 'items/day', percentiles, growth_rate)
            PowerPointGenerator.add_two_column_slide(prs, "Items per Day", items_bullets, items_bytes)
        else:
            PowerPointGenerator.add_two_column_slide(prs, "Items per Day",
                ["Quantity column not provided.", "Items/day chart unavailable.",
                 "Map a numeric quantity field in 'Column mapping' if available."])
        
        if bin_summary is None:
            bin_bullets = ["No bin estimation generated.", "Run 'Estimate Bins' to populate this slide.",
                          "This section based on manual assumptions + order profile velocity."]
            PowerPointGenerator.add_two_column_slide(prs, "Bin Dimensioning & Stock-on-hand", bin_bullets)
        else:
            bl, bw, bh = bin_summary.get("bin_dims", (None, None, None))
            peak_months = bin_summary.get("peak_months")
            peak_txt = ", ".join([datetime(2000, m, 1).strftime("%b") for m in peak_months]) if peak_months else "(not applied)"
            
            bin_bullets = [
                f"Estimated total bins: {format_integer(bin_summary.get('total_bins'))} (across {format_integer(bin_summary.get('skus_counted'))} SKUs).",
                f"Design basis: {bin_summary.get('planning_basis', 'N/A')}.",
                f"Stock-on-hand: {format_integer(bin_summary.get('stock_on_hand_days'))} days.",
                f"Safety stock: {format_percentage(bin_summary.get('safety_stock_pct'))}.",
                f"Bin utilisation: {format_percentage(bin_summary.get('bin_util_pct'))}.",
                f"Items per bin: {format_number(bin_summary.get('items_per_bin'), 0)}.",
                f"Bin dims (L/W/H): {bl:g} / {bw:g} / {bh:g}.", f"Peak months filter: {peak_txt}.",
                "Note: estimate sensitive to rounding at SKU level."
            ]
            
            if bin_df is not None and not bin_df.empty and "ABC_Class" in bin_df.columns:
                try:
                    by_abc = bin_df.groupby("ABC_Class")["Bins"].sum().sort_values(ascending=False)
                    top = [f"{idx}: {format_integer(val)} bins" for idx, val in by_abc.head(3).items()]
                    if top: bin_bullets.append("Bins by ABC (top): " + " | ".join(top))
                except: pass
            
            PowerPointGenerator.add_two_column_slide(prs, "Bin Dimensioning & Stock-on-hand", bin_bullets)
        
        if batch_summary is None or batch_results is None:
            batch_bullets = ["No batching analysis performed.", "Run 'Calculate Batch Factor' to populate.",
                           "Batching groups orders by similarity to optimize picking."]
            PowerPointGenerator.add_two_column_slide(prs, "Order Batching Analysis", batch_bullets)
        else:
            batch_fig = ChartGenerator.create_batch_factor_chart(batch_results, figsize=(5.2, 5.2), dpi=170)
            batch_bytes = PowerPointGenerator.figure_to_bytes(batch_fig)
            plt.close(batch_fig)
            
            batch_bullets = [
                f"Total groups: {format_integer(batch_summary.get('total_groups'))}",
                f"Group size: {batch_summary.get('group_size')} orders/batch",
                f"Waving timer: {batch_summary.get('waving_timer')} min" + (" (all day)" if batch_summary.get('waving_timer') == 0 else ""),
                f"Overall avg batch factor: {batch_summary.get('overall_avg', 0):.4f}",
                f"Overall median batch factor: {batch_summary.get('overall_median', 0):.4f}",
                f"Similarity threshold: {batch_summary.get('threshold')}",
                f"MinHash permutations: {batch_summary.get('permutations')}",
                "Batch Factor = Total items / Unique items",
                "Higher values = better batching efficiency (more item overlap)."
            ]
            PowerPointGenerator.add_two_column_slide(prs, "Order Batching Analysis", batch_bullets, batch_bytes)
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()


# ============================================================================
# SESSION STATE
# ============================================================================

class SessionState:
    @staticmethod
    def initialize():
        defaults = {
            'analysis_ran': False, 'abc_df': None, 'abc_summary': None, 'daily_df': None,
            'daily_summary': None, 'series_choice': 'Lines', 'yoy_growth': Config.DEFAULT_GROWTH,
            'pptx_bytes': None, 'pptx_key': None, 'bin_df': None, 'bin_summary': None,
            'batch_df': None, 'batch_results': None, 'batch_summary': None,
            'csv_encoding': "Auto", 'csv_delimiter': "Auto",
        }
        for key, value in defaults.items():
            if key not in st.session_state:
                st.session_state[key] = value


# ============================================================================
# UI COMPONENTS
# ============================================================================

def render_page_config():
    st.set_page_config(page_title=Config.PAGE_TITLE, layout=Config.PAGE_LAYOUT)
    st.markdown("""
        <style>
            .block-container {max-width: 1180px; padding-top: 1.25rem; padding-bottom: 2.0rem;}
            h1, h2, h3 {margin-bottom: 0.35rem; letter-spacing: -0.2px;}
            .stCaption {color: rgba(17,17,17,0.65);}
            .stMarkdown p {margin-bottom: 0.25rem;}
            [data-testid="stDataFrame"] {border-radius: 12px; overflow: hidden; border: 1px solid rgba(0,0,0,0.06);}
            .stButton>button {padding: 0.55rem 0.95rem; border-radius: 12px; font-weight: 600;}
            [data-testid="stMetric"] {border-radius: 14px; border: 1px solid rgba(0,0,0,0.06); padding: 0.35rem 0.25rem;}
            hr {margin: 1.0rem 0;}
        </style>
    """, unsafe_allow_html=True)

def render_sidebar(uploaded_file):
    with st.sidebar:
        st.header("Inputs")
        st.file_uploader("Upload file (.xlsx or .csv)", type=['xlsx', 'csv'], key='uploaded')
        
        with st.expander("CSV import options", expanded=False):
            st.selectbox("Encoding", options=Config.CSV_ENCODING_OPTIONS, key="csv_encoding",
                        help="Auto detects. Override if needed.")
            st.selectbox("Delimiter", options=Config.CSV_DELIMITER_OPTIONS, key="csv_delimiter",
                        help="Auto sniffs. Override if needed.")
            st.caption("Multi-pass fallback: C engine → Python → skip bad lines → replacement.")
        
        st.divider()
        st.subheader("ABC thresholds")
        a_cutoff = st.slider("A cutoff", 0.50, 0.95, Config.DEFAULT_A_CUTOFF, 0.01)
        b_cutoff = st.slider("B cutoff", 0.70, 0.99, Config.DEFAULT_B_CUTOFF, 0.01)
        if b_cutoff <= a_cutoff:
            st.error("B cutoff must be greater than A cutoff.")
        
        st.divider()
        st.subheader("Percentiles (daily)")
        percentiles = st.multiselect("Show percentile lines",
            [0.50, 0.55, 0.60, 0.65, 0.70, 0.75, 0.80, 0.90, 0.95, 1.00],
            default=Config.DEFAULT_PERCENTILES, format_func=lambda x: f"P{int(x*100)}")
    
    return {'a_cutoff': a_cutoff, 'b_cutoff': b_cutoff, 'percentiles': percentiles, 'valid': b_cutoff > a_cutoff}

def render_summary_metrics(abc_summary, daily_summary, has_quantity):
    st.divider()
    st.subheader("Summary")
    cols = st.columns(4, gap="medium")
    cols[0].metric("Total SKUs", format_integer(abc_summary.get('total_skus')))
    cols[1].metric("Total Items", format_number(abc_summary.get('total_value'), 0))
    if has_quantity:
        cols[2].metric("Peak Items/day", format_integer(daily_summary.get('peak_qty_value')))
    else:
        cols[2].metric("A share", format_percentage(abc_summary.get('A_value_share')))
    cols[3].metric("Peak Lines/day", format_integer(daily_summary.get('peak_lines_value')))

def render_abc_analysis(abc_df, abc_summary, a_cutoff, b_cutoff):
    st.subheader("ABC Distribution")
    left_col, right_col = st.columns([2, 1], gap="large")
    
    with left_col:
        fig, _, _ = ChartGenerator.create_abc_chart(abc_df, abc_summary, a_cutoff, b_cutoff)
        st.pyplot(fig, use_container_width=False)
        plt.close(fig)
    
    with right_col:
        st.caption("SKU concentration (Lines %)")
        concentration_data = []
        for pct in [0.10, 0.20, 0.30]:
            if 'SKU_Pct' in abc_df.columns:
                total_skus = abc_summary.get('total_skus', len(abc_df))
                target_sku = max(1, min(int(round(pct * total_skus)), len(abc_df)))
                mask = abc_df['SKU_Count'] == target_sku
                idx = abc_df[mask].index[0] if mask.any() else (abc_df['SKU_Pct'] - pct).abs().idxmin()
                lines_share = float(abc_df.loc[idx, 'cum_share']) * 100.0
            else:
                n = len(abc_df)
                k = min(max(1, int(round(pct * max(abc_summary.get('total_skus', n), n)))), n)
                lines_share = float(abc_df.iloc[k - 1]['cum_share']) * 100.0
            concentration_data.append({'Top SKUs': f"Top {int(pct*100)}%", 'Lines %': f"{lines_share:.1f}%"})
        st.table(pd.DataFrame(concentration_data))

def render_daily_analysis(daily_df, percentiles):
    st.subheader("Lines / Items per Day")
    left_col, right_col = st.columns([2, 1], gap="large")
    has_quantity = 'Quantity' in daily_df.columns and daily_df['Quantity'].notna().any()
    
    with left_col:
        options = ['Lines'] if not has_quantity else ['Lines', 'Items']
        st.radio("Series", options=options, horizontal=True, key='series_choice')
        st.number_input("Total Growth Projection (%)", -100.0, 300.0, float(st.session_state.yoy_growth),
                       1.0, key='yoy_growth', help="Annual growth rate for scenario projection.")
        
        y_col = 'Lines' if st.session_state.series_choice == 'Lines' else 'Quantity'
        y_label = 'Lines/day' if st.session_state.series_choice == 'Lines' else 'Items/day'
        
        fig = ChartGenerator.create_daily_volume_chart(daily_df, y_col, y_label, percentiles, st.session_state.yoy_growth)
        st.pyplot(fig, use_container_width=False)
        plt.close(fig)
    
    with right_col:
        st.caption("Peak + percentiles (scenario basis)")
        plot_df = daily_df[['Date', y_col]].copy()
        plot_df['Date'] = pd.to_datetime(plot_df['Date'], errors='coerce', format="%Y-%m-%d")
        if plot_df['Date'].isna().all():
            plot_df['Date'] = pd.to_datetime(plot_df['Date'], errors='coerce')
        plot_df[y_col] = pd.to_numeric(plot_df[y_col], errors='coerce')
        plot_df = plot_df.dropna(subset=['Date', y_col]).sort_values('Date')
        
        growth_factor = 1.0 + (st.session_state.yoy_growth / 100.0)
        plot_df['Projected'] = plot_df[y_col] * growth_factor
        basis_col = 'Projected' if st.session_state.yoy_growth != 0.0 else y_col
        
        if not plot_df.empty:
            peak_row = plot_df.loc[plot_df[basis_col].idxmax()]
            peak_label = 'Peak lines' if y_col == 'Lines' else 'Peak items'
            st.write("Peak day:", peak_row['Date'].date().isoformat())
            st.write(f"{peak_label}:", format_integer(peak_row[basis_col]))
            values = plot_df[basis_col].to_numpy()
            values = values[np.isfinite(values)]
            for pct in sorted(percentiles):
                st.write(f"P{int(pct*100)}:", format_integer(np.quantile(values, pct)))

def detect_abc_class_column(abc_df):
    return first_existing_col(abc_df, ['ABC', 'Class', 'ABC_Class', 'abc_class', 'abc']) if abc_df is not None else None

def compute_daily_lines_per_sku(df, sku_col, date_col):
    tmp = df[[sku_col, date_col]].copy()
    tmp['Date'] = safe_to_datetime_date(tmp[date_col])
    tmp = tmp.dropna(subset=[sku_col, 'Date'])
    daily = tmp.groupby([sku_col, 'Date']).size().reset_index(name='Lines')
    daily['Lines'] = pd.to_numeric(daily['Lines'], errors='coerce').fillna(0.0)
    return daily

def compute_bin_estimate(daily_sku_df, sku_col, planning_basis, stock_on_hand_days, safety_stock_pct,
                         items_per_bin, bin_util_pct, peak_months=None):
    work = daily_sku_df.copy()
    work['Date'] = pd.to_datetime(work['Date'], errors='coerce')
    work = work.dropna(subset=['Date', sku_col])
    work['Lines'] = pd.to_numeric(work['Lines'], errors='coerce').fillna(0.0)
    
    if peak_months:
        work = work[work['Date'].dt.month.isin(peak_months)].copy()
    if work.empty:
        return pd.DataFrame(columns=[sku_col, 'Velocity', 'StockOnHand', 'Bins'])
    
    basis = planning_basis.strip().upper()
    if basis == "AVG":
        vel = work.groupby(sku_col)['Lines'].mean()
    else:
        q_map = {"P80": 0.80, "P90": 0.90, "P95": 0.95, "P99": 0.99}
        vel = work.groupby(sku_col)['Lines'].quantile(q_map.get(basis, 0.95))
    
    vel = vel.to_frame(name='Velocity').reset_index()
    util = clamp01(bin_util_pct)
    eff_capacity = max(1e-9, float(items_per_bin) * max(1e-9, util))
    
    vel['StockOnHand'] = vel['Velocity'] * int(stock_on_hand_days) * (1.0 + float(safety_stock_pct))
    vel['Bins'] = np.ceil(vel['StockOnHand'] / eff_capacity).astype(int)
    vel['Velocity'], vel['StockOnHand'] = vel['Velocity'].astype(float), vel['StockOnHand'].astype(float)
    return vel.sort_values('Bins', ascending=False).reset_index(drop=True)

def render_bin_estimator(df, sku_col, date_col, abc_df):
    st.subheader("Bin Count Estimator (no dims/weights)")
    left, right = st.columns([2, 1], gap="large")
    
    with left:
        st.caption("Assumptions + planning basis")
        planning_basis = st.selectbox("Planning Basis (daily lines per SKU)",
            ["AVG", "P80", "P90", "P95", "P99"],
            index=["AVG", "P80", "P90", "P95", "P99"].index(Config.DEFAULT_PLANNING_BASIS),
            help="Percentile of daily lines per SKU as velocity.")
        
        use_peak_months = st.checkbox("Use peak months only", value=False)
        peak_months = None
        if use_peak_months:
            peak_months = st.multiselect("Peak months", list(range(1, 13)), default=[11, 12],
                                        format_func=lambda m: datetime(2000, m, 1).strftime("%b"))
        
        st.divider()
        st.caption("Inventory policy (proxy)")
        stock_on_hand_days = st.number_input("Stock-on-hand (days)", 1, 365, int(Config.DEFAULT_STOCK_ON_HAND_DAYS), 1)
        safety_stock_pct = st.number_input("Empty bins (%)", 0.0, 300.0, float(Config.DEFAULT_SAFETY_STOCK * 100.0), 1.0) / 100.0
        
        st.divider()
        st.caption("Bin assumptions")
        bin_util_pct = st.number_input("Bin utilisation (%)", 1.0, 100.0, float(Config.DEFAULT_BIN_UTIL * 100.0), 1.0) / 100.0
        items_per_bin = st.number_input("Items per bin", 1.0, 1_000_000.0, float(Config.DEFAULT_ITEMS_PER_BIN), 1.0,
                                       help="Capacity assumption when no dims/weights exist.")
        
        st.divider()
        st.caption("Bin dimensions")
        c1, c2, c3 = st.columns(3)
        bin_l = c1.number_input("Length (mm)", 0.0, value=float(Config.DEFAULT_BIN_L), step=1.0)
        bin_w = c2.number_input("Width (mm)", 0.0, value=float(Config.DEFAULT_BIN_W), step=1.0)
        bin_h = c3.number_input("Height (mm)", 0.0, value=float(Config.DEFAULT_BIN_H), step=1.0)
        st.caption("Dims stored for traceability. Don't alter math unless you map size→capacity.")
    
    abc_class_col = detect_abc_class_column(abc_df)
    abc_map = None
    if abc_df is not None and abc_class_col and sku_col in abc_df.columns:
        abc_map = abc_df[[sku_col, abc_class_col]].dropna().drop_duplicates(subset=[sku_col]).rename(columns={abc_class_col: "ABC_Class"})
    
    with right:
        st.caption("Outputs")
        if st.button("Estimate Bins", use_container_width=True):
            try:
                daily_sku_df = compute_daily_lines_per_sku(df, sku_col, date_col)
                bin_df = compute_bin_estimate(daily_sku_df, sku_col, planning_basis, int(stock_on_hand_days),
                                             float(safety_stock_pct), float(items_per_bin), float(bin_util_pct),
                                             peak_months if use_peak_months else None)
                
                if not bin_df.empty:
                    bin_df["Bin_L"], bin_df["Bin_W"], bin_df["Bin_H"] = float(bin_l), float(bin_w), float(bin_h)
                    bin_df["Bin_Util"], bin_df["ItemsPerBin"] = float(bin_util_pct), float(items_per_bin)
                    bin_df["SOH_Days"], bin_df["SafetyStock"] = int(stock_on_hand_days), float(safety_stock_pct)
                    bin_df["PlanningBasis"] = planning_basis
                
                if abc_map is not None and not bin_df.empty:
                    bin_df = bin_df.merge(abc_map, on=sku_col, how='left')
                elif 'ABC_Class' not in bin_df.columns:
                    bin_df['ABC_Class'] = np.nan
                
                total_bins = int(bin_df['Bins'].sum()) if not bin_df.empty else 0
                sku_count = int(bin_df[sku_col].nunique()) if not bin_df.empty else 0
                
                st.session_state.bin_df = bin_df
                st.session_state.bin_summary = {
                    "planning_basis": planning_basis, "stock_on_hand_days": int(stock_on_hand_days),
                    "safety_stock_pct": float(safety_stock_pct), "bin_util_pct": float(bin_util_pct),
                    "items_per_bin": float(items_per_bin), "bin_dims": (float(bin_l), float(bin_w), float(bin_h)),
                    "total_bins": total_bins, "skus_counted": sku_count,
                    "peak_months": peak_months if use_peak_months else None, "line_mode": "rows",
                }
            except Exception as e:
                st.error(f"Bin estimate failed: {e}")
                st.session_state.bin_df = st.session_state.bin_summary = None
        
        if st.session_state.get('bin_df') is not None and st.session_state.get('bin_summary') is not None:
            bsum, bdf = st.session_state.bin_summary, st.session_state.bin_df
            st.metric("Estimated total bins", format_integer(bsum["total_bins"]))
            st.write("SKUs counted:", format_integer(bsum["skus_counted"]))
            st.write("Planning basis:", bsum["planning_basis"])
            st.write("Stock-on-hand (days):", bsum["stock_on_hand_days"])
            st.write("Safety stock:", format_percentage(bsum["safety_stock_pct"]))
            st.write("Bin utilisation:", format_percentage(bsum["bin_util_pct"]))
            st.write("Items per bin:", format_number(bsum["items_per_bin"], 0))
            bl, bw, bh = bsum["bin_dims"]
            st.write("Bin dims (L/W/H):", f"{bl:g} / {bw:g} / {bh:g}")
            if bsum["peak_months"]:
                st.write("Peak months:", ", ".join([datetime(2000, m, 1).strftime("%b") for m in bsum["peak_months"]]))
            
            if "ABC_Class" in bdf.columns and bdf["ABC_Class"].notna().any():
                by_abc = bdf.groupby("ABC_Class")["Bins"].sum().sort_values(ascending=False).reset_index()
                st.divider()
                st.caption("Bins by ABC class")
                st.dataframe(by_abc, use_container_width=True, hide_index=True)

def render_batching_calculator(df, order_col, sku_col, date_col):
    st.subheader("Order Batching Calculator")
    left, right = st.columns([2, 1], gap="large")
    
    with left:
        st.caption("Batching parameters")
        waving_timer = st.number_input("Waving Timer (minutes, 0 = all day)", 0, 1440, Config.DEFAULT_WAVING_TIMER, 5,
                                      help="Time interval for grouping orders. 0 = group all orders by date only.")
        group_size = st.number_input("Group Size (orders per batch)", 2, 50, Config.DEFAULT_GROUP_SIZE, 1,
                                    help="Number of orders to group together for picking.")
        
        st.divider()
        st.caption("Similarity settings")
        threshold = st.number_input("Jaccard Threshold", 0.0, 1.0, Config.DEFAULT_THRESHOLD, 0.01, format="%.2f",
                                   help="Minimum similarity (0-1) for grouping. 0.1 = 10% overlap.")
        permutations = st.number_input("MinHash Permutations", 32, 512, Config.DEFAULT_PERMUTATIONS, 32,
                                      help="Number of hash functions. Higher = more accurate but slower.")
    
    with right:
        st.caption("Calculate batching")
        if st.button("Calculate Batch Factor", use_container_width=True):
            try:
                with st.spinner("Processing batching analysis..."):
                    batch_df = df[[date_col, order_col, sku_col]].copy()
                    
                    if pd.api.types.is_numeric_dtype(batch_df[date_col]):
                        start_date = pd.to_datetime("2024-01-01")
                        batch_df[date_col] = batch_df[date_col].astype(int).apply(lambda x: start_date + pd.Timedelta(days=x - 1))
                    else:
                        batch_df[date_col] = pd.to_datetime(batch_df[date_col].astype(str), errors="coerce", dayfirst=True)
                    
                    batch_df = batch_df.dropna(subset=[date_col])
                    
                    if waving_timer > 0:
                        if not batch_df[date_col].dt.time.ne(pd.to_datetime("00:00:00").time()).any():
                            st.error("Waving Timer > 0 but no time data in date column.")
                            return
                        batch_df["roundeddatetime"] = batch_df[date_col].apply(lambda x: round_up_datetime(x, waving_timer))
                    else:
                        batch_df["roundeddatetime"] = batch_df[date_col].dt.date
                    
                    batch_df["roundeddatetime"] = pd.to_datetime(batch_df["roundeddatetime"], errors="coerce")
                    batch_df = batch_df.dropna(subset=["roundeddatetime", order_col, sku_col])
                    
                    grouped_results = group_clients_by_day_optimized(batch_df, int(group_size), int(permutations),
                                                                     float(threshold), order_col, sku_col)
                    
                    out_df = pd.DataFrame(grouped_results)
                    out_df['Order'] = out_df['Order'].apply(lambda x: ";".join(map(str, x)))
                    out_df.sort_values(["roundeddatetime", "Taskgroup"], inplace=True)
                    
                    df_report = out_df[out_df["Taskgroup"].notna()].copy()
                    batch_factors = df_report["Batch Factor"].tolist()
                    overall_avg = sum(batch_factors) / len(batch_factors) if batch_factors else 0
                    overall_median = statistics.median(batch_factors) if batch_factors else 0
                    
                    st.session_state.batch_df = out_df
                    st.session_state.batch_results = grouped_results
                    st.session_state.batch_summary = {
                        "total_groups": len(grouped_results), "group_size": int(group_size),
                        "waving_timer": int(waving_timer), "threshold": float(threshold),
                        "permutations": int(permutations), "overall_avg": float(overall_avg),
                        "overall_median": float(overall_median),
                    }
                    st.success(f"✓ Formed {len(grouped_results)} batch groups")
            except Exception as e:
                st.error(f"Batching calculation failed: {e}")
                import traceback
                st.code(traceback.format_exc())
        
        if st.session_state.get('batch_summary') is not None:
            bsum = st.session_state.batch_summary
            st.divider()
            st.metric("Total Groups", format_integer(bsum["total_groups"]))
            st.write("Avg Batch Factor:", f"{bsum['overall_avg']:.4f}")
            st.write("Median Batch Factor:", f"{bsum['overall_median']:.4f}")
            st.write("Group Size:", format_integer(bsum["group_size"]))
            st.write("Waving Timer:", f"{bsum['waving_timer']} min" if bsum['waving_timer'] > 0 else "All day")
    
    if st.session_state.get('batch_results') is not None:
        st.divider()
        fig = ChartGenerator.create_batch_factor_chart(st.session_state.batch_results)
        st.pyplot(fig, use_container_width=True)
        plt.close(fig)


# ============================================================================
# MAIN APPLICATION
# ============================================================================

def main():
    render_page_config()
    SessionState.initialize()
    
    st.title(Config.PAGE_TITLE)
    st.caption("Upload .xlsx or .csv. Run ABC, Daily Volumes, Bin Estimation, and Order Batching analysis.")
    
    uploaded = st.session_state.get('uploaded')
    config = render_sidebar(uploaded)
    
    if uploaded is None:
        st.info("Upload a .xlsx or .csv file in the sidebar to begin.")
        st.stop()
    
    if not config['valid']:
        st.stop()
    
    try:
        file_bytes = uploaded.getvalue()
        filename = uploaded.name.lower()
        
        if filename.endswith('.csv'):
            df, message = DataLoader.load_csv(file_bytes, uploaded.name,
                st.session_state.get("csv_encoding", "Auto"), st.session_state.get("csv_delimiter", "Auto"))
            st.caption(message)
            sheet = None
            show_preview = st.toggle("Show data preview", value=False)
        elif filename.endswith('.xlsx'):
            xls = pd.ExcelFile(uploaded)
            col1, col2 = st.columns([2, 1], gap="large")
            with col1: sheet = st.selectbox("Sheet", xls.sheet_names)
            with col2: show_preview = st.toggle("Show data preview", value=False)
            df = DataLoader.load_excel(xls, sheet_name=sheet)
            st.caption(f"Loaded Excel: {uploaded.name} (Sheet: {sheet})")
        else:
            st.error("Unsupported file type. Please upload .xlsx or .csv.")
            st.stop()
        
        if show_preview:
            with st.expander("Data preview (first 20 rows)", expanded=False):
                st.dataframe(df.head(20), use_container_width=True)
    except Exception as e:
        st.error(f"Could not load file: {e}")
        st.stop()
    
    if df is None or df.empty:
        st.error("No data found in file.")
        st.stop()
    
    columns = df.columns.tolist()
    if not columns:
        st.error("No columns detected in file.")
        st.stop()
    
    st.subheader("Column mapping")
    col1, col2 = st.columns(2, gap="large")
    with col1:
        sku_col = st.selectbox("SKU No.", columns)
        line_col = st.selectbox("Line No.", columns)
        order_col = st.selectbox("Order ID", columns, help="Order identifier for batching analysis")
    with col2:
        date_col = st.selectbox("Date Input", columns)
        qty_option = st.selectbox("Quantity No.", ["(none)"] + columns)
        qty_col = None if qty_option == "(none)" else qty_option
    
    if st.button("Run Analysis", use_container_width=True):
        try:
            abc_df, abc_summary = run_abc(df, sku_col, line_col, a_cut=config['a_cutoff'], b_cut=config['b_cutoff'], debug=False)
            daily_df, daily_summary = run_lines_per_day(df, date_col, qty_col, tuple(sorted(config['percentiles'])))
            
            if daily_df is not None and not daily_df.empty and "Date" in daily_df.columns:
                daily_df["Date"] = coerce_to_yyyy_mm_dd(daily_df["Date"])
            
            st.session_state.abc_df, st.session_state.abc_summary = abc_df, abc_summary
            st.session_state.daily_df, st.session_state.daily_summary = daily_df, daily_summary
            st.session_state.analysis_ran = True
            st.session_state.pptx_bytes = st.session_state.pptx_key = None
            st.session_state.bin_df = st.session_state.bin_summary = None
            st.session_state.batch_df = st.session_state.batch_results = st.session_state.batch_summary = None
            
            has_qty = 'Quantity' in daily_df.columns and daily_df['Quantity'].notna().any()
            if not has_qty and st.session_state.series_choice == 'Items':
                st.session_state.series_choice = 'Lines'
        except Exception as e:
            st.error(f"Analysis failed: {e}")
            st.session_state.analysis_ran = False
    
    if not st.session_state.analysis_ran:
        st.info("Set inputs, then click **Run Analysis**.")
        st.stop()
    
    abc_df, abc_summary = st.session_state.abc_df, st.session_state.abc_summary
    daily_df, daily_summary = st.session_state.daily_df, st.session_state.daily_summary
    has_quantity = 'Quantity' in daily_df.columns and daily_df['Quantity'].notna().any()
    
    render_summary_metrics(abc_summary, daily_summary, has_quantity)
    render_abc_analysis(abc_df, abc_summary, config['a_cutoff'], config['b_cutoff'])
    render_daily_analysis(daily_df, config['percentiles'])
    
    st.divider()
    render_bin_estimator(df, sku_col, date_col, abc_df)
    
    st.divider()
    render_batching_calculator(df, order_col, sku_col, date_col)
    
    with st.expander("Open results tables", expanded=False):
        st.markdown("**ABC Distribution Table**")
        st.dataframe(abc_df, use_container_width=True)
        st.download_button("Download ABC Distribution (CSV)", df_to_csv_bytes(abc_df), "abc_distribution.csv",
                          "text/csv", use_container_width=True)
        
        st.markdown("**Daily Volumes Table**")
        st.dataframe(daily_df, use_container_width=True)
        st.download_button("Download Daily Volumes (CSV)", df_to_csv_bytes(daily_df), "daily_volumes.csv",
                          "text/csv", use_container_width=True)
        
        st.markdown("**Bins Estimator Table**")
        if st.session_state.get('bin_df') is not None:
            bin_df = st.session_state.bin_df
            st.dataframe(bin_df, use_container_width=True)
            st.download_button("Download Bins Estimator (CSV)", df_to_csv_bytes(bin_df), "bins_estimator.csv",
                              "text/csv", use_container_width=True)
        else:
            st.caption("No bin estimate available. Click 'Estimate Bins' above.")
        
        st.markdown("**Batching Results Table**")
        if st.session_state.get('batch_df') is not None:
            display_df = st.session_state.batch_df.copy()
            if 'unique_items' in display_df.columns:
                display_df['unique_items'] = display_df['unique_items'].apply(
                    lambda x: ';'.join(str(i) for i in x) if isinstance(x, set) else x)
            st.dataframe(display_df, use_container_width=True)
            st.download_button("Download Batching Results (CSV)", df_to_csv_bytes(display_df), "batching_results.csv",
                              "text/csv", use_container_width=True)
        else:
            st.caption("No batching results available. Click 'Calculate Batch Factor' above.")
    
    st.divider()
    st.subheader("Report")
    
    def get_or_build_report():
        cache_key = (
            uploaded.name, sheet, sku_col, line_col, date_col, qty_col,
            float(config['a_cutoff']), float(config['b_cutoff']), tuple(sorted(config['percentiles'])),
            float(st.session_state.yoy_growth),
            None if st.session_state.get('bin_summary') is None else (
                st.session_state.bin_summary.get("planning_basis"), st.session_state.bin_summary.get("stock_on_hand_days"),
                st.session_state.bin_summary.get("safety_stock_pct"), st.session_state.bin_summary.get("bin_util_pct"),
                st.session_state.bin_summary.get("items_per_bin"), st.session_state.bin_summary.get("bin_dims"),
                tuple(st.session_state.bin_summary.get("peak_months") or []), st.session_state.bin_summary.get("total_bins"),
                st.session_state.bin_summary.get("skus_counted"),
            ),
            None if st.session_state.get('batch_summary') is None else (
                st.session_state.batch_summary.get("total_groups"), st.session_state.batch_summary.get("group_size"),
                st.session_state.batch_summary.get("waving_timer"), st.session_state.batch_summary.get("threshold"),
                st.session_state.batch_summary.get("permutations"),
            )
        )
        
        if st.session_state.pptx_bytes is None or st.session_state.pptx_key != cache_key:
            st.session_state.pptx_bytes = PowerPointGenerator.generate_report(
                uploaded.name, sheet, sku_col, line_col, date_col, qty_col,
                float(config['a_cutoff']), float(config['b_cutoff']), list(sorted(config['percentiles'])),
                float(st.session_state.yoy_growth), abc_df, abc_summary, daily_df, df,
                st.session_state.get('bin_summary'), st.session_state.get('bin_df'),
                st.session_state.get('batch_summary'), st.session_state.get('batch_results'),
            )
            st.session_state.pptx_key = cache_key
        return st.session_state.pptx_bytes
    
    report_bytes = get_or_build_report()
    st.download_button("Generate Report (PowerPoint)", report_bytes, "enhanced_analysis_report.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)


if __name__ == "__main__":
    main()